# -*- coding: utf-8 -*-
import sys
import csv
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.text import MSO_ANCHOR

def set_line_color(line, rgb_color):
    """Establece el color de la línea de forma compatible con varias versiones."""
    try:
        line.color.rgb = rgb_color
    except AttributeError:
        try:
            line.fill.solid()
            line.fill.fore_color.rgb = rgb_color
        except AttributeError:
            print(f"Advertencia: No se pudo establecer el color de la línea a {rgb_color}")

def generar_presentacion(csv_file, pptx_file):
    """
    Lee datos de un fichero CSV y genera una presentación de PowerPoint
    con hasta 8 cartas por slide distribuidas uniformemente con un recuadro por cada valor de campo en el centro (ajuste automático) y conversión a string.

    Args:
        csv_file (str): Ruta al fichero CSV de entrada.
        pptx_file (str): Ruta al fichero PPTX de salida.
    """
    try:
        with open(csv_file, 'r', encoding='utf-8') as f:
            reader = csv.DictReader(f)
            if reader.fieldnames is None:
                print(f"Error: El fichero CSV '{csv_file}' está vacío o no tiene encabezado.")
                return

            prs = Presentation()
            slide_layout = prs.slide_layouts[6]  # Usamos el layout completamente en blanco

            # Dimensiones de la carta y márgenes para distribución uniforme
            slide_ancho_util_in = prs.slide_width - Inches(1) # Margen de 0.5 pulgadas a cada lado
            slide_alto_util_in = prs.slide_height - Inches(1) # Margen de 0.5 pulgadas arriba y abajo
            cartas_por_fila = 4
            filas_por_slide = 2
            espacio_horizontal_total = slide_ancho_util_in
            espacio_vertical_total = slide_alto_util_in
            carta_ancho_max_in = espacio_horizontal_total / cartas_por_fila * 0.9 # Dejar un 10% para márgenes
            carta_alto_max_in = espacio_vertical_total / filas_por_slide * 0.9 # Dejar un 10% para márgenes
            carta_aspecto = 63 / 88 # Relación de aspecto de la carta original

            if carta_ancho_max_in / carta_aspecto > carta_alto_max_in:
                carta_alto_in = carta_alto_max_in
                carta_ancho_in = carta_alto_in * carta_aspecto
            else:
                carta_ancho_in = carta_ancho_max_in
                carta_alto_in = carta_ancho_in / carta_aspecto

            margen_horizontal = (slide_ancho_util_in - cartas_por_fila * carta_ancho_in) / (cartas_por_fila + 1)
            margen_vertical = (slide_alto_util_in - filas_por_slide * carta_alto_in) / (filas_por_slide + 1)

            # Dimensiones proporcionales de los elementos dentro de la carta
            titulo_alto_in = carta_alto_in * 0.15
            titulo_margen_top_in = carta_alto_in * 0.02
            habilidad_alto_in = carta_alto_in * 0.2
            habilidad_margen_bottom_in = carta_alto_in * 0.02
            texto_margen_interno = carta_ancho_in * 0.05 # Margen interno para el texto central
            espacio_entre_campos = carta_alto_in * 0.01 # Espacio vertical entre recuadros de campos
            fuente_base_texto = 8 # Tamaño de fuente base para el texto central

            carta_index = 0
            slide = None
            cartas_en_slide = 0

            for row in reader:
                if carta_index % 8 == 0:
                    slide = prs.slides.add_slide(slide_layout)
                    cartas_en_slide = 0

                fila = cartas_en_slide // cartas_por_fila
                columna = cartas_en_slide % cartas_por_fila

                left = Inches(0.5) + (columna + 1) * margen_horizontal + columna * carta_ancho_in
                top = Inches(0.5) + (fila + 1) * margen_vertical + fila * carta_alto_in

                # Área de la carta
                carta_forma = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    left,
                    top,
                    carta_ancho_in,
                    carta_alto_in
                )
                fill = carta_forma.fill
                fill.solid()
                fill.fore_color.rgb = RGBColor(220, 220, 220)
                line = carta_forma.line
                set_line_color(line, RGBColor(0, 0, 0))

                # Recuadro del título
                titulo_rect_forma = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    left + carta_ancho_in * 0.02,
                    top + titulo_margen_top_in,
                    carta_ancho_in * 0.96,
                    titulo_alto_in
                )
                titulo_rect_forma.fill.solid()
                titulo_rect_forma.fill.fore_color.rgb = RGBColor(0, 0, 139)
                titulo_rect_line = titulo_rect_forma.line
                set_line_color(titulo_rect_line, RGBColor(0, 0, 0))
                titulo_frame = titulo_rect_forma.text_frame
                titulo_frame.clear()
                titulo_paragraph = titulo_frame.paragraphs[0]
                titulo_text = row.get("Nombre del Perfil", "Título").upper()
                titulo_paragraph.text = titulo_text
                titulo_paragraph.alignment = PP_ALIGN.CENTER
                titulo_frame.vertical_anchor = 1 # MSO_ANCHOR.MIDDLE
                titulo_run = titulo_paragraph.runs[0]
                titulo_run.font.bold = True
                titulo_run.font.color.rgb = RGBColor(255, 165, 0)
                titulo_run.font.size = Pt(int(10 * carta_alto_in / carta_alto_max_in)) # Ajuste de tamaño

                # Área para los recuadros de los campos centrales
                texto_area_left = left + texto_margen_interno
                texto_area_top = top + titulo_alto_in + carta_alto_in * 0.02
                texto_area_width = carta_ancho_in - 2 * texto_margen_interno
                texto_area_height = top + carta_alto_in - habilidad_alto_in - habilidad_margen_bottom_in - texto_area_top

                campos_centrales = {key: value for key, value in row.items() if key not in ["Nombre del Perfil", "Habilidad Especial", "Mazo"]}
                num_campos = len(campos_centrales)

                if num_campos > 0:
                    campo_height = (texto_area_height - (num_campos - 1) * espacio_entre_campos) / num_campos
                    campo_top_actual = texto_area_top
                    font_size_pt = int(fuente_base_texto * carta_alto_in / carta_alto_max_in)

                    for key, value in campos_centrales.items():
                        campo_rect = slide.shapes.add_textbox(
                            texto_area_left,
                            campo_top_actual,
                            texto_area_width,
                            campo_height
                        )
                        text_frame = campo_rect.text_frame
                        text_frame.clear()
                        p = text_frame.paragraphs[0]
                        run = p.add_run()
                        run.text = str(value) # Convertir el valor a string
                        font = run.font
                        font.size = Pt(font_size_pt)
                        text_frame.vertical_anchor = MSO_ANCHOR.TOP # Alineación vertical al top
                        p.alignment = PP_ALIGN.LEFT # Justificar a la izquierda

                        # Añadir la clave en una línea separada y más pequeña encima del valor
                        clave_p = text_frame.add_paragraph()
                        clave_run = clave_p.add_run()
                        clave_run.text = f"{key}:"
                        clave_run.font.size = Pt(font_size_pt * 0.8) # Tamaño más pequeño para la clave
                        clave_p.alignment = PP_ALIGN.LEFT

                        campo_top_actual += campo_height + espacio_entre_campos

                # Recuadro de la habilidad especial (texto centrado)
                habilidad_rect_forma = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    left + carta_ancho_in * 0.02,
                    top + carta_alto_in - habilidad_alto_in - habilidad_margen_bottom_in,
                    carta_ancho_in * 0.96,
                    habilidad_alto_in
                )
                habilidad_rect_forma.fill.solid()
                habilidad_rect_forma.fill.fore_color.rgb = RGBColor(255, 215, 0)
                habilidad_rect_line = habilidad_rect_forma.line
                set_line_color(habilidad_rect_line, RGBColor(0, 0, 0))
                habilidad_frame = habilidad_rect_forma.text_frame
                habilidad_frame.clear()
                habilidad_paragraph = habilidad_frame.paragraphs[0]
                habilidad_text = row.get("Habilidad Especial", "").upper()
                habilidad_paragraph.text = habilidad_text
                habilidad_paragraph.alignment = PP_ALIGN.CENTER # Centrar horizontalmente
                habilidad_frame.vertical_anchor = MSO_ANCHOR.MIDDLE # Centrar verticalmente
                habilidad_run = habilidad_paragraph.runs[0]
                habilidad_run.font.size = Pt(int(7 * carta_alto_in / carta_alto_max_in)) # Ajuste de tamaño
                habilidad_run.font.italic = True
                habilidad_run.font.color.rgb = RGBColor(0, 0, 0)

                carta_index += 1
                cartas_en_slide += 1

            prs.save(pptx_file)
            print(f"Se ha generado el fichero PowerPoint: {pptx_file}")

    except FileNotFoundError:
        print(f"Error: No se encontró el fichero CSV '{csv_file}'.")
    except Exception as e:
        print(f"Ocurrió un error: {e}")

if __name__ == "__main__":
    csv_file_name = "cartas.csv"
    pptx_file_name = "cartas_tablero_string_conversion.pptx" # Nuevo nombre de archivo
    if len(sys.argv) > 1:
        csv_file_name = sys.argv[1]
    if len(sys.argv) > 2:
        pptx_file_name = sys.argv[2]
    generar_presentacion(csv_file_name, pptx_file_name)
